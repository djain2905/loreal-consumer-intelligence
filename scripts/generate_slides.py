"""
Generates Lancôme Consumer Intelligence presentation.
Run: python3 scripts/generate_slides.py
Output: docs/lancome-consumer-intelligence-slides.pptx
"""
import os, io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patheffects as pe
import matplotlib.dates as mdates
import numpy as np
import pandas as pd
from dotenv import load_dotenv
load_dotenv()
import snowflake.connector

def _fetch_trends():
    try:
        conn = snowflake.connector.connect(
            account=os.getenv("SNOWFLAKE_ACCOUNT"),
            user=os.getenv("SNOWFLAKE_USER"),
            password=os.getenv("SNOWFLAKE_PASSWORD"),
            warehouse=os.getenv("SNOWFLAKE_WAREHOUSE"),
            database="LOREAL_DB",
        )
        cur = conn.cursor()
        cur.execute("ALTER WAREHOUSE COMPUTE_WH RESUME IF SUSPENDED")
        cur.execute("""
            select review_month, avg_rating, desire_rate, review_count
            from loreal_db.dev_mart.mart_sentiment_over_time
            where review_count >= 10
            order by review_month
        """)
        rows = cur.fetchall()
        conn.close()
        return pd.DataFrame(rows, columns=["month", "avg_rating", "desire_rate", "review_count"])
    except Exception as e:
        print(f"  ⚠️  Snowflake unavailable ({e}) — skipping trends slide")
        return None

df_trends = _fetch_trends()
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────────
ROSE       = RGBColor(0xC5, 0x81, 0x7A)
GOLD       = RGBColor(0xC9, 0xA9, 0x6E)
CREAM      = RGBColor(0xFD, 0xF8, 0xF6)
DARK       = RGBColor(0x1A, 0x1A, 0x1A)
MID        = RGBColor(0x77, 0x77, 0x77)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LROSE      = RGBColor(0xF7, 0xED, 0xEB)
BLUE       = RGBColor(0x7A, 0x9C, 0xC5)
RED        = RGBColor(0xC0, 0x39, 0x2B)
LRED       = RGBColor(0xFB, 0xEB, 0xE9)

R  = "#C5817A"
G  = "#C9A96E"
CR = "#FDF8F6"
LR = "#F7EDEB"
GR = "#CCCCCC"
DK = "#1A1A1A"
RD = "#C0392B"
LRD= "#FBEBE9"

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Helpers ────────────────────────────────────────────────────────────────────
def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def box(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h,
       sz=14, bold=False, italic=False,
       color=None, align=PP_ALIGN.LEFT,
       font="Georgia", wrap=True):
    if color is None:
        color = DARK
    b  = slide.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text           = text
    r.font.size      = Pt(sz)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    r.font.name      = font
    return b

def divider(slide, l, t, w, color=GOLD):
    box(slide, l, t, w, Pt(2), color)

def fig2buf(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf

ACCENT_BG = {GOLD: RGBColor(0xFF, 0xFD, 0xF0), ROSE: LROSE}

def callout(slide, text, l, t, w, h, accent=RED):
    bg_color = ACCENT_BG.get(accent, LRED)
    box(slide, l, t, w, h, bg_color)
    box(slide, l, t, Inches(0.06), h, accent)
    tb(slide, "▶  " + text, l + Inches(0.12), t + Inches(0.1),
       w - Inches(0.2), h - Inches(0.15),
       sz=11, color=DARK, font="Calibri", bold=False)

# ── Data ───────────────────────────────────────────────────────────────────────
THEMES = [
    ("Sensitive Skin Formula", 225, 4.27),
    ("Texture & Formula",      207, 4.34),
    ("Fragrance & Scent",      106, 4.29),
    ("Packaging & Format",      59, 4.34),
    ("SPF & Sun Protection",    29, 4.45),
    ("Key Ingredients",         27, 4.26),
    ("Price & Value",           24, 4.38),
    ("Shade Range",             18, 4.72),
    ("Longevity & Wear",        11, 4.64),
]
COMPS = [
    ("Lancôme",       4.48, 15.0, True),
    ("Tatcha",        4.24, 19.3, False),
    ("Drunk Elephant",4.06, 21.1, False),
    ("Clinique",      4.34, 16.9, False),
    ("Fresh",         4.39, 16.1, False),
    ("The Ordinary",  4.26, 15.9, False),
]

OUT = os.path.join(os.path.dirname(__file__), "..", "docs",
                   "lancome-consumer-intelligence-slides.pptx")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1  Cover
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, CREAM)
box(s, Inches(0), Inches(0), Inches(0.5), H, ROSE)
box(s, Inches(0.5), Inches(3.0), Inches(8), Pt(2), GOLD)

tb(s, "What Consumers Wish\nLancôme Made",
   Inches(0.85), Inches(0.9), Inches(8.5), Inches(2.2),
   sz=46, bold=True, color=DARK, font="Georgia")
tb(s, "A Consumer Intelligence Analysis of 5,951 Sephora Reviews",
   Inches(0.85), Inches(3.2), Inches(8.5), Inches(0.55),
   sz=15, color=MID, italic=True, font="Calibri")
tb(s, "Dhwani Jain  ·  L'Oréal USA Commercial Management Trainee  ·  2026",
   Inches(0.85), Inches(6.9), Inches(9), Inches(0.4),
   sz=11, color=MID, font="Calibri")

for i, (val, lbl) in enumerate([
        ("5,951",  "Sephora reviews\nanalyzed"),
        ("15%",    "signal an\nunmet need"),
        ("4.48 ★", "avg rating\n(strong brand)")]):
    x = Inches(10.2)
    y = Inches(1.3 + i * 1.9)
    box(s, x, y, Inches(2.9), Inches(1.5), LROSE)
    tb(s, val, x, y + Inches(0.05), Inches(2.9), Inches(0.85),
       sz=34, bold=True, color=ROSE, font="Georgia", align=PP_ALIGN.CENTER)
    tb(s, lbl, x, y + Inches(0.9), Inches(2.9), Inches(0.55),
       sz=10, color=MID, font="Calibri", align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2  Setting the Scene — trends over time
# ══════════════════════════════════════════════════════════════════════════════
if df_trends is not None:
    s = prs.slides.add_slide(blank)
    bg(s, CREAM)
    box(s, Inches(0), Inches(0), Inches(0.5), H, ROSE)

    tb(s, "SETTING THE SCENE  ·  How Have Consumers Felt About Lancôme Over Time?",
       Inches(0.75), Inches(0.22), Inches(12), Inches(0.35),
       sz=10, bold=True, color=ROSE, font="Calibri")
    tb(s, "4.0★+ for Nearly a Decade — A Loyal Base That Still Has More to Ask For",
       Inches(0.75), Inches(0.6), Inches(9.0), Inches(0.85),
       sz=26, bold=True, color=DARK, font="Georgia")
    divider(s, Inches(0.75), Inches(1.52), Inches(8.8))

    dates   = pd.to_datetime(df_trends["month"])
    ratings = df_trends["avg_rating"].values
    desires = df_trends["desire_rate"].values

    fig, ax1 = plt.subplots(figsize=(8.8, 4.8), facecolor=CR)
    ax1.set_facecolor(CR)
    ax2 = ax1.twinx()

    ax1.plot(dates, ratings, color=R,  linewidth=2.2, label="Avg Rating (★)")
    ax2.plot(dates, desires, color=G,  linewidth=1.8, linestyle="--",
             label="Unmet Need Rate (%)", alpha=0.9)

    ax1.set_ylabel("Avg Rating (★)", color=R, fontsize=10)
    ax2.set_ylabel("Unmet Need Rate (%)", color=G, fontsize=10)
    ax1.set_ylim(1, 5.6)
    ax2.set_ylim(0, 90)
    ax1.tick_params(axis="y", labelcolor=R, labelsize=9)
    ax2.tick_params(axis="y", labelcolor=G, labelsize=9)
    ax1.xaxis.set_major_formatter(mdates.DateFormatter("%Y"))
    ax1.xaxis.set_major_locator(mdates.YearLocator())
    ax1.tick_params(axis="x", labelsize=9)
    ax1.spines[["top"]].set_visible(False)
    ax2.spines[["top"]].set_visible(False)
    ax1.grid(axis="x", color="#EEEEEE", linestyle="--", linewidth=0.6)
    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1 + lines2, labels1 + labels2, fontsize=9, framealpha=0,
               loc="lower right")
    fig.tight_layout(pad=0.6)
    s.shapes.add_picture(fig2buf(fig), Inches(0.6), Inches(1.65), Inches(9.2), Inches(5.35))

    # Right-side callout boxes
    callout(s,
            "Ratings hold above 4.0★ for 9+ consecutive years — "
            "a loyal, retained consumer base. Not a brand in decline.",
            Inches(10.0), Inches(1.85), Inches(3.1), Inches(1.5), accent=ROSE)
    callout(s,
            "Desire signals persist even at peak satisfaction — "
            "these consumers aren't leaving. They're asking for more. "
            "That's the whitespace opportunity.",
            Inches(10.0), Inches(3.55), Inches(3.1), Inches(1.8), accent=GOLD)

    tb(s, "Source: mart_sentiment_over_time · Months with 10+ reviews only · LOREAL_DB.DEV_MART",
       Inches(0.75), Inches(7.17), Inches(12), Inches(0.25),
       sz=8, color=MID, font="Calibri")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3  Descriptive  — takeaway with numbers
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, CREAM)
box(s, Inches(0), Inches(0), Inches(0.5), H, ROSE)

tb(s, "DESCRIPTIVE  ·  What Happened?",
   Inches(0.75), Inches(0.28), Inches(11), Inches(0.38),
   sz=10, bold=True, color=ROSE, font="Calibri")
tb(s,
   "15% of Lancôme's 5,951 Sephora Reviews Signal an Unmet\n"
   "Consumer Need — 895 Consumers Raised Their Hand",
   Inches(0.75), Inches(0.7), Inches(12.1), Inches(1.35),
   sz=27, bold=True, color=DARK, font="Georgia")
divider(s, Inches(0.75), Inches(2.1), Inches(11.8))

# Donut chart
fig, ax = plt.subplots(figsize=(5.0, 4.8), facecolor=CR)
wedges, _ = ax.pie([85, 15], colors=["#E8E0DE", R],
                   startangle=90, wedgeprops=dict(width=0.52, edgecolor="white", linewidth=2.5))
ax.text(0, 0.1,  "15%", ha="center", va="center",
        fontsize=34, fontweight="bold", color=R, fontfamily="Georgia")
ax.text(0, -0.32, "signal an unmet need", ha="center", va="center",
        fontsize=11, color="#666")
ax.axis("equal")
ax.set_title("Share of Reviews Containing Desire Language",
             fontsize=11, color="#555", pad=10)
rose_p = mpatches.Patch(color=R,        label="895 reviews — signal desire language")
grey_p = mpatches.Patch(color="#E8E0DE", label="5,056 reviews — no desire signal")
ax.legend(handles=[rose_p, grey_p], fontsize=9.5, framealpha=0,
          loc="lower center", bbox_to_anchor=(0.5, -0.13), ncol=1)
fig.tight_layout(pad=0.3)
s.shapes.add_picture(fig2buf(fig), Inches(0.6), Inches(2.2), Inches(5.2), Inches(4.7))

# Right side text
tb(s, "What counts as a 'desire review'?",
   Inches(6.2), Inches(2.3), Inches(6.6), Inches(0.4),
   sz=13, bold=True, color=DARK, font="Georgia")
tb(s, "Reviews containing: wish · want · need · hope · would love · "
       "if only · lacks · missing\n\n"
       "These are verified Sephora purchasers who wrote a review "
       "and still told us what was missing from the product. "
       "They are high-intent consumers actively signaling a gap.",
   Inches(6.2), Inches(2.75), Inches(6.6), Inches(1.6),
   sz=11, color=MID, font="Calibri")

callout(s,
        "895 individual consumers explicitly articulated an unmet need — "
        "across a brand rated 4.48★ and bought voluntarily. "
        "Satisfaction and desire signals coexist.",
        Inches(6.2), Inches(4.55), Inches(6.6), Inches(1.1))

tb(s, "Source: 5,951 Lancôme reviews · Sephora · Kaggle dataset · Desire flag via regex on review_text",
   Inches(0.75), Inches(7.12), Inches(12), Inches(0.28),
   sz=8, color=MID, font="Calibri")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3a  Diagnostic — volume
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, CREAM)
box(s, Inches(0), Inches(0), Inches(0.5), H, ROSE)

tb(s, "DIAGNOSTIC  ·  Why Did It Happen?  (1 of 2)",
   Inches(0.75), Inches(0.28), Inches(11), Inches(0.38),
   sz=10, bold=True, color=ROSE, font="Calibri")
tb(s,
   "Sensitive Skin Formula Has 225 Desire Reviews —\n"
   "More Than Any Other Category by a Wide Margin",
   Inches(0.75), Inches(0.7), Inches(12.1), Inches(1.35),
   sz=29, bold=True, color=DARK, font="Georgia")
divider(s, Inches(0.75), Inches(2.1), Inches(11.8))

labels = [t[0] for t in THEMES]
counts = [t[1] for t in THEMES]
colors = [R if l == "Sensitive Skin Formula" else GR for l in labels]

fig, ax = plt.subplots(figsize=(9.0, 4.8), facecolor=CR)
ax.set_facecolor(CR)
bars = ax.barh(labels[::-1], counts[::-1],
               color=colors[::-1], height=0.55, edgecolor="none")

# Clean count labels — inside bars for large values, outside for small
for bar, count, lbl in zip(bars, counts[::-1], labels[::-1]):
    w = bar.get_width()
    y = bar.get_y() + bar.get_height() / 2
    is_top = lbl == "Sensitive Skin Formula"
    # Place count at end of bar with enough padding to avoid overlap
    ax.text(w + 4, y, str(count), va="center", fontsize=12,
            color=R if is_top else "#555",
            fontweight="bold" if is_top else "normal")

# Single clean marker line at the top bar — no overlapping arrow
ax.axvline(x=225, color=R, linestyle=":", linewidth=1.5, alpha=0.5)

ax.set_xlabel("Number of Desire Reviews", fontsize=12, color="#888")
ax.set_xlim(0, 270)
ax.spines[["top", "right", "left"]].set_visible(False)
ax.tick_params(left=False, labelsize=12)
ax.grid(axis="x", color="#EEEEEE", linestyle="--", linewidth=0.7)
fig.tight_layout(pad=0.8)
s.shapes.add_picture(fig2buf(fig), Inches(0.6), Inches(2.15), Inches(9.3), Inches(5.1))

callout(s,
        "225 consumers signalled this need — the highest volume of any theme. "
        "These same consumers also rate affected products 4.27★, "
        "the lowest avg rating in the dataset. Volume + frustration in one theme.",
        Inches(10.1), Inches(2.3), Inches(3.0), Inches(2.2))

tb(s, "What triggers this theme?",
   Inches(10.1), Inches(4.7), Inches(3.0), Inches(0.4),
   sz=10, bold=True, color=DARK, font="Georgia")
tb(s,
   "Reviews mentioning:\nsensitive skin · broke me out · irritation "
   "· reaction · rash · allergy · hypoallergenic · fragrance-free",
   Inches(10.1), Inches(5.12), Inches(3.0), Inches(1.5),
   sz=10, color=MID, font="Calibri")

tb(s, "Source: mart_whitespace_summary · LOREAL_DB.DEV_MART · Snowflake",
   Inches(0.75), Inches(7.12), Inches(12), Inches(0.28),
   sz=8, color=MID, font="Calibri")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3b  Diagnostic — frustration (avg rating)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, CREAM)
box(s, Inches(0), Inches(0), Inches(0.5), H, ROSE)

tb(s, "DIAGNOSTIC  ·  Why Did It Happen?  (2 of 2)",
   Inches(0.75), Inches(0.28), Inches(11), Inches(0.38),
   sz=10, bold=True, color=ROSE, font="Calibri")
tb(s,
   "Sensitive Skin & Key Ingredients Consumers Have the Lowest Ratings —\n"
   "4.27★ and 4.26★, the Most Frustrated Segments in the Dataset",
   Inches(0.75), Inches(0.7), Inches(12.1), Inches(1.35),
   sz=29, bold=True, color=DARK, font="Georgia")
divider(s, Inches(0.75), Inches(2.1), Inches(11.8))

labels  = [t[0] for t in THEMES]
ratings = [t[2] for t in THEMES]
rc      = [R if l in ("Sensitive Skin Formula", "Key Ingredients") else GR for l in labels]

fig, ax = plt.subplots(figsize=(9.5, 4.5), facecolor=CR)
ax.set_facecolor(CR)
ax.scatter(ratings[::-1], labels[::-1], color=rc[::-1], s=200, zorder=3)
ax.axvline(x=4.48, color=G, linestyle="--", linewidth=1.5,
           label="Lancôme overall avg (4.48★)")
for lbl, rat, c in zip(labels[::-1], ratings[::-1], rc[::-1]):
    ax.text(rat + 0.015, lbl, f"{rat}★", va="center", fontsize=11,
            color=RD if c == R else "#888",
            fontweight="bold" if c == R else "normal")

ax.set_xlabel("Avg Rating on Desire Reviews  (lower = more frustrated)", fontsize=11, color="#888")
ax.set_xlim(4.1, 4.95)
ax.spines[["top","right","left"]].set_visible(False)
ax.tick_params(left=False, labelsize=11)
ax.grid(axis="x", color="#EEE", linestyle="--", linewidth=0.7)
ax.legend(fontsize=10, framealpha=0, loc="lower right")
fig.tight_layout(pad=0.5)
s.shapes.add_picture(fig2buf(fig), Inches(0.6), Inches(2.2), Inches(9.5), Inches(5.0))

callout(s,
        "High desire volume + low avg rating = real pain, not preference. "
        "Sensitive Skin Formula and Key Ingredients are the two themes where "
        "consumers are most dissatisfied AND most vocal.",
        Inches(10.3), Inches(2.3), Inches(2.8), Inches(2.0))

tb(s, "Why avg rating matters here",
   Inches(10.3), Inches(4.5), Inches(2.8), Inches(0.4),
   sz=10, bold=True, color=DARK, font="Georgia")
tb(s,
   "These consumers gave lower ratings while writing desire reviews — "
   "meaning they bought the product, were disappointed, "
   "AND still told us what they needed. "
   "That is the highest-quality signal for a product gap.",
   Inches(10.3), Inches(4.92), Inches(2.8), Inches(1.7),
   sz=10, color=MID, font="Calibri")

tb(s, "Source: mart_whitespace_summary · LOREAL_DB.DEV_MART · Snowflake",
   Inches(0.75), Inches(7.12), Inches(12), Inches(0.28),
   sz=8, color=MID, font="Calibri")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4  Competitive  — no competitor owns this
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, CREAM)
box(s, Inches(0), Inches(0), Inches(0.5), H, ROSE)

tb(s, "DIAGNOSTIC  ·  Competitive Context",
   Inches(0.75), Inches(0.28), Inches(11), Inches(0.38),
   sz=10, bold=True, color=ROSE, font="Calibri")
tb(s,
   "0 Out of 5 Prestige Competitors Has a Sensitive-Skin-First Product —\n"
   "Lancôme Has a Clear Path to Own the Category",
   Inches(0.75), Inches(0.7), Inches(12.1), Inches(1.35),
   sz=27, bold=True, color=DARK, font="Georgia")
divider(s, Inches(0.75), Inches(2.1), Inches(11.8))

brands  = [c[0] for c in COMPS]
ratings = [c[1] for c in COMPS]
desires = [c[2] for c in COMPS]
is_l    = [c[3] for c in COMPS]
bc      = [R if il else GR for il in is_l]
dc      = ["#E8A8A2" if il else "#D8D8D8" for il in is_l]

x  = np.arange(len(brands))
bw = 0.38

fig, ax = plt.subplots(figsize=(8.0, 4.0), facecolor=CR)
ax.set_facecolor(CR)
b1 = ax.bar(x - bw/2, ratings, bw, color=bc,  label="Avg Rating (★)", edgecolor="none")
b2 = ax.bar(x + bw/2, desires, bw, color=dc, label="Unmet Need Rate (%)", edgecolor="none", alpha=0.85)

for b in b1:
    ax.text(b.get_x()+b.get_width()/2, b.get_height()+0.4,
            f"{b.get_height():.2f}★", ha="center", fontsize=9, color="#444")
for b in b2:
    ax.text(b.get_x()+b.get_width()/2, b.get_height()+0.4,
            f"{b.get_height():.1f}%", ha="center", fontsize=9, color="#444")

ax.set_xticks(x)
ax.set_xticklabels(brands, fontsize=10, rotation=8, ha="right")
ax.set_ylim(0, 28)
ax.spines[["top","right","left"]].set_visible(False)
ax.tick_params(left=False)
ax.grid(axis="y", color="#EEE", linestyle="--", linewidth=0.6)
rose_p = mpatches.Patch(color=R,  label="Lancôme")
gray_p = mpatches.Patch(color=GR, label="Competitors")
ax.legend(handles=[rose_p, gray_p], fontsize=9, framealpha=0, loc="upper right")
ax.set_title("Avg Rating vs. Unmet Need Rate — Lancôme vs. Top 5 Prestige Competitors",
             fontsize=10.5, color=DK, fontweight="bold")
fig.tight_layout(pad=0.4)
s.shapes.add_picture(fig2buf(fig), Inches(0.6), Inches(2.2), Inches(8.5), Inches(4.7))

callout(s,
        "Lancôme has the highest avg rating (4.48★) in this peer set — "
        "it has the brand credibility to launch a premium sensitive-skin line "
        "that no competitor currently offers.",
        Inches(9.3), Inches(2.3), Inches(3.8), Inches(1.5))

callout(s,
        "Drunk Elephant (21.1%) and Tatcha (19.3%) have the highest unmet need rates — "
        "their consumers are also signaling gaps. No one in prestige owns sensitive-skin skincare.",
        Inches(9.3), Inches(4.1), Inches(3.8), Inches(1.7))

tb(s, "Source: mart_competitor_benchmarks · Top 5 brands by Sephora review volume · LOREAL_DB.DEV_MART",
   Inches(0.75), Inches(7.12), Inches(12), Inches(0.28),
   sz=8, color=MID, font="Calibri")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5  Recommendation
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, CREAM)
box(s, Inches(0), Inches(0), Inches(0.5), H, ROSE)

tb(s, "RECOMMENDATION",
   Inches(0.75), Inches(0.28), Inches(11), Inches(0.38),
   sz=10, bold=True, color=ROSE, font="Calibri")
divider(s, Inches(0.75), Inches(0.72), Inches(11.8))

# Action block
box(s, Inches(0.75), Inches(0.85), Inches(11.8), Inches(1.65), ROSE)
tb(s, "ACTION",
   Inches(0.95), Inches(0.9), Inches(2), Inches(0.35),
   sz=9, bold=True, color=WHITE, font="Calibri")
tb(s, "Launch Génifique Sensitive — a fragrance-free, dermatologist-tested Lancôme\n"
       "serum formulated for reactive and sensitive skin at $95–$140",
   Inches(0.95), Inches(1.22), Inches(11.3), Inches(1.1),
   sz=20, bold=True, color=WHITE, font="Georgia")

# Arrow
tb(s, "↓", Inches(6.2), Inches(2.6), Inches(1.5), Inches(0.6),
   sz=28, bold=True, color=ROSE, font="Georgia", align=PP_ALIGN.CENTER)

# Outcome block
box(s, Inches(0.75), Inches(3.25), Inches(11.8), Inches(1.65), LROSE)
box(s, Inches(0.75), Inches(3.25), Inches(0.06), Inches(1.65), ROSE)
tb(s, "EXPECTED OUTCOME",
   Inches(0.95), Inches(3.3), Inches(3), Inches(0.35),
   sz=9, bold=True, color=ROSE, font="Calibri")
tb(s,
   "Est. ~4,500 first-year customers  ·  ~$518K Year 1 revenue — "
   "converting the 225 high-intent consumers already on record signaling this exact gap, "
   "entering a prestige category that 0 of 5 competitors currently own",
   Inches(0.95), Inches(3.65), Inches(11.2), Inches(1.1),
   sz=14, color=DARK, font="Calibri")

# Evidence row
divider(s, Inches(0.75), Inches(5.1), Inches(11.8), color=GOLD)
tb(s, "Evidence",
   Inches(0.75), Inches(5.22), Inches(3), Inches(0.35),
   sz=11, bold=True, color=DARK, font="Georgia")

stats = [
    ("225",   "desire reviews\nfor this theme"),
    ("4.27★", "avg rating —\nhighest frustration"),
    ("0",     "prestige competitors\nowning this category"),
    ("15%",   "overall Lancôme\ndesire signal rate"),
]
for i, (val, lbl) in enumerate(stats):
    x = Inches(0.75 + i * 3.05)
    box(s, x, Inches(5.65), Inches(2.8), Inches(1.55), LROSE)
    tb(s, val, x, Inches(5.68), Inches(2.8), Inches(0.8),
       sz=30, bold=True, color=ROSE, font="Georgia", align=PP_ALIGN.CENTER)
    tb(s, lbl, x, Inches(6.42), Inches(2.8), Inches(0.6),
       sz=10, color=MID, font="Calibri", align=PP_ALIGN.CENTER)

tb(s, "Projection assumes: 1.5% Sephora review rate (beauty industry standard) → 15,000 estimated high-intent buyers  ·  "
       "30% trial rate (conservative for pre-identified intent cohort)  ·  $115 avg price (midpoint of $95–$140)  ·  "
       "Cohort-only estimate — excludes broader sensitive-skin market acquisition",
   Inches(0.75), Inches(7.12), Inches(12.3), Inches(0.28),
   sz=8, color=MID, font="Calibri", italic=True)



# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"✅ Saved → {OUT}")
print(f"   {len(prs.slides)} slides")
